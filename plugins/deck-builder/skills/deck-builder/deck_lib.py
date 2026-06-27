"""Opinionated python-pptx wrapper for the deck-builder skill.

The model calls NAMED layout methods (title_slide, content_slide, …) and never hand-rolls
geometry — that keeps decks consistent and on-brand. `python-pptx` is imported lazily
inside `Deck` so the pure helpers (PALETTE, hex_to_rgb, image_fit, fit_text) import with no
third-party dependency and stay unit-testable.

Customize by editing PALETTE / FOOTER, or pass your own palette to Deck(palette=...).
"""
from __future__ import annotations

# --- Brand (override freely) -------------------------------------------------
PALETTE = {
    "primary": "2E5AAC",       # headers, accents
    "primary_dark": "1B355F",  # title/section backgrounds
    "text": "1A1A1A",          # body text
    "muted": "6B7280",         # captions, footer
    "surface": "F3F4F6",       # card / tile fills
    "white": "FFFFFF",
    "accent": "12B5A5",        # stat highlights
}
FOOTER = ""  # set a footer label (e.g. your name/org) or leave empty


# --- Pure helpers (no python-pptx needed) ------------------------------------
def hex_to_rgb(h: str) -> tuple[int, int, int]:
    """'2E5AAC' or '#2E5AAC' -> (46, 90, 172)."""
    h = h.lstrip("#")
    if len(h) != 6 or any(c not in "0123456789abcdefABCDEF" for c in h):
        raise ValueError(f"invalid hex color: {h!r}")
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def image_fit(img_w: float, img_h: float, box_w: float, box_h: float) -> tuple[float, float]:
    """Scale (img_w, img_h) to fit inside (box_w, box_h) preserving aspect ratio.

    Returns the rendered (w, h). Setting both from this — rather than forcing width AND
    height independently — is what prevents distortion (a real failure mode when decks are
    imported into Google Slides)."""
    if img_w <= 0 or img_h <= 0:
        raise ValueError("image dimensions must be positive")
    scale = min(box_w / img_w, box_h / img_h)
    return img_w * scale, img_h * scale


def fit_text(text: str, max_chars: int) -> str:
    """Single-line truncation with an ellipsis, for headlines/labels that must not wrap."""
    text = " ".join(text.split())
    if max_chars <= 1 or len(text) <= max_chars:
        return text
    return text[: max_chars - 1].rstrip() + "…"


# --- The deck (lazy python-pptx) ---------------------------------------------
class Deck:
    """Thin builder over python-pptx. 16:9. Footer + page number auto-applied.

    Usage:
        d = Deck("My Deck")
        d.title_slide("My Deck", "A subtitle")
        d.content_slide("Batch genealogy needs brittle manual joins", ["point", "point"])
        d.stat_tiles("By the numbers", [("3x", "faster"), ("0", "manual steps")])
        d.section_divider("Approach")
        d.closing_slide("Thank you", "you@example.com")
        d.save("out.pptx")
    """

    def __init__(self, title: str = "Deck", *, palette: dict | None = None, footer: str | None = None):
        from pptx import Presentation  # lazy: only needed to actually build
        from pptx.util import Inches

        self.palette = {**PALETTE, **(palette or {})}
        self.footer = FOOTER if footer is None else footer
        self._title = title
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self._blank = self.prs.slide_layouts[6]  # blank layout
        self._n = 0

    # -- internal helpers -----------------------------------------------------
    def _rgb(self, key: str):
        from pptx.dml.color import RGBColor
        return RGBColor(*hex_to_rgb(self.palette[key]))

    def _bg(self, slide, key: str):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self._rgb(key)

    def _add_slide(self, bg: str = "white"):
        slide = self.prs.slides.add_slide(self._blank)
        self._bg(slide, bg)
        self._n += 1
        self._chrome(slide)
        return slide

    def _textbox(self, slide, left, top, width, height, text, *, size, color, bold=False, align="left"):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = self._rgb(color)
        return box

    def _chrome(self, slide):
        # Footer label + page number in muted text.
        if self.footer:
            self._textbox(slide, 0.4, 7.05, 8, 0.3, self.footer, size=9, color="muted")
        self._textbox(slide, 12.4, 7.05, 0.6, 0.3, str(self._n), size=9, color="muted", align="right")

    # -- named layouts --------------------------------------------------------
    def title_slide(self, title: str, subtitle: str = ""):
        slide = self._add_slide("primary_dark")
        self._textbox(slide, 0.8, 2.6, 11.7, 1.5, title, size=40, color="white", bold=True)
        if subtitle:
            self._textbox(slide, 0.8, 4.0, 11.7, 1.0, subtitle, size=20, color="surface")
        return slide

    def section_divider(self, title: str):
        slide = self._add_slide("primary")
        self._textbox(slide, 0.8, 3.1, 11.7, 1.3, title, size=32, color="white", bold=True)
        return slide

    def content_slide(self, headline: str, bullets: list[str]):
        """Headline should be a full assertion (the takeaway), not a topic label."""
        slide = self._add_slide("white")
        self._textbox(slide, 0.8, 0.6, 11.7, 1.2, headline, size=26, color="primary_dark", bold=True)
        from pptx.util import Inches, Pt
        box = slide.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.6))
        tf = box.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(10)
            run = p.add_run()
            run.text = f"•  {b}"
            run.font.size = Pt(18)
            run.font.color.rgb = self._rgb("text")
        return slide

    def stat_tiles(self, headline: str, tiles: list[tuple[str, str]]):
        """tiles: list of (big_value, caption). Up to 4 across."""
        slide = self._add_slide("white")
        self._textbox(slide, 0.8, 0.6, 11.7, 1.0, headline, size=26, color="primary_dark", bold=True)
        from pptx.util import Inches
        n = max(1, min(len(tiles), 4))
        gap, margin = 0.4, 0.9
        tile_w = (13.333 - 2 * margin - (n - 1) * gap) / n
        for i, (value, caption) in enumerate(tiles[:n]):
            left = margin + i * (tile_w + gap)
            card = slide.shapes.add_shape(1, Inches(left), Inches(2.3), Inches(tile_w), Inches(2.4))
            card.fill.solid(); card.fill.fore_color.rgb = self._rgb("surface")
            card.line.color.rgb = self._rgb("surface")
            self._textbox(slide, left, 2.6, tile_w, 1.1, value, size=40, color="accent", bold=True, align="center")
            self._textbox(slide, left, 3.7, tile_w, 0.8, caption, size=14, color="muted", align="center")
        return slide

    def closing_slide(self, title: str, contact: str = ""):
        slide = self._add_slide("primary_dark")
        self._textbox(slide, 0.8, 3.0, 11.7, 1.2, title, size=36, color="white", bold=True)
        if contact:
            self._textbox(slide, 0.8, 4.2, 11.7, 0.7, contact, size=18, color="surface")
        return slide

    def notes(self, slide, text: str):
        """Attach speaker notes — the talk track, not a reprint of the slide."""
        slide.notes_slide.notes_text_frame.text = text

    def save(self, path: str):
        self.prs.save(path)
        return path
